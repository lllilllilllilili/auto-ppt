from __future__ import annotations

import base64
import os
from io import BytesIO

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn

from app.models import SlideRequest, ImageData

TEMPLATE_PATH = os.path.join(os.path.dirname(__file__), "assets", "template.pptx")

LAYOUT_TITLE = 0      # 제목 슬라이드
LAYOUT_CONTENT = 1    # 제목 및 내용
LAYOUT_COMPARE = 3    # 비교

COVER_BG = RGBColor(0x00, 0x1D, 0x2C)


def _set_slide_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _remove_placeholder(slide, idx):
    """placeholder를 슬라이드에서 완전히 제거"""
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            sp = ph._element
            sp.getparent().remove(sp)
            return


def _set_placeholder_text(ph, text, font_name="맑은 고딕", font_size=None,
                          font_color=None, bold=None, alignment=None):
    tf = ph.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    if font_name:
        p.font.name = font_name
    if font_size:
        p.font.size = font_size
    if font_color:
        p.font.color.rgb = font_color
    if bold is not None:
        p.font.bold = bold
    if alignment is not None:
        p.alignment = alignment


def _add_textbox(slide, left, top, width, height, text, font_name="맑은 고딕",
                 font_size=Pt(14), font_color=RGBColor(0x33, 0x33, 0x33),
                 bold=False, alignment=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True

    anchor_map = {MSO_ANCHOR.TOP: "t", MSO_ANCHOR.MIDDLE: "ctr", MSO_ANCHOR.BOTTOM: "b"}
    body_pr = tf._txBody.find(qn("a:bodyPr"))
    if body_pr is not None:
        body_pr.set("anchor", anchor_map.get(anchor, "t"))

    p = tf.paragraphs[0]
    p.text = text
    p.font.name = font_name
    p.font.size = font_size
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.alignment = alignment
    return txbox


def _fill_bullets_in_placeholder(ph, bullets, font_name="맑은 고딕", font_size=Pt(14)):
    tf = ph.text_frame
    tf.clear()
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if i > 0:
            p.space_before = Pt(6)
        p.text = b
        p.font.name = font_name
        p.font.size = font_size


def _add_bullets(slide, left, top, width, height, bullets,
                 font_name="맑은 고딕", font_size=Pt(14),
                 text_color=RGBColor(0x33, 0x33, 0x33)):
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True

    for i, bullet_text in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if i > 0:
            p.space_before = Pt(6)
        p.text = bullet_text
        p.font.name = font_name
        p.font.size = font_size
        p.font.color.rgb = text_color

        pPr = p._p.get_or_add_pPr()
        for tag in ["a:buChar", "a:buFont", "a:buClr"]:
            for old in pPr.findall(qn(tag)):
                pPr.remove(old)
        pPr.append(pPr.makeelement(qn("a:buFont"), {"typeface": "Arial"}))
        pPr.append(pPr.makeelement(qn("a:buChar"), {"char": "●"}))
        pPr.set("indent", str(Emu(-228600)))
        pPr.set("marL", str(Emu(457200)))

    return txbox


def _add_divider(slide, left, top, width, color, thickness=Pt(1)):
    shape = slide.shapes.add_shape(1, left, top, width, thickness)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _decode_image(img_data: ImageData) -> BytesIO:
    raw = img_data.data
    if "," in raw:
        raw = raw.split(",", 1)[1]
    return BytesIO(base64.b64decode(raw))


def _get_image_size(stream: BytesIO):
    """PIL로 이미지 원본 크기를 읽고 stream을 되감음"""
    from PIL import Image as PILImage
    img = PILImage.open(stream)
    w, h = img.size
    stream.seek(0)
    return w, h


def _add_image(slide, img_data: ImageData, left, top, max_width, max_height):
    """원본 비율을 유지하면서 max_width x max_height 안에 맞춤"""
    stream = _decode_image(img_data)
    orig_w, orig_h = _get_image_size(stream)

    aspect = orig_w / orig_h
    box_aspect = max_width / max_height

    if aspect > box_aspect:
        width = max_width
        height = int(max_width / aspect)
    else:
        height = max_height
        width = int(max_height * aspect)

    slide.shapes.add_picture(stream, left, top, width, height)


# 슬라이드 영역: 제목바 아래 콘텐츠 영역 (10in x ~4.5in)
IMAGE_POSITIONS = {
    "left-top":    {"left": 0.3,  "top": 1.0, "w": 4.0, "h": 2.8},
    "right-top":   {"left": 5.5,  "top": 1.0, "w": 4.2, "h": 2.8},
    "center":      {"left": 1.5,  "top": 1.0, "w": 7.0, "h": 3.5},
    "full":        {"left": 0.3,  "top": 1.0, "w": 9.4, "h": 4.2},
    "left-bottom": {"left": 0.3,  "top": 3.0, "w": 4.0, "h": 2.2},
    "right-bottom":{"left": 5.5,  "top": 3.0, "w": 4.2, "h": 2.2},
    "left-half":   {"left": 0.3,  "top": 1.0, "w": 4.5, "h": 4.2},
}


def _add_positioned_image(slide, img_data: ImageData):
    pos = IMAGE_POSITIONS.get(img_data.position, IMAGE_POSITIONS["right-top"])
    _add_image(slide, img_data,
               Inches(pos["left"]), Inches(pos["top"]),
               Inches(pos["w"]), Inches(pos["h"]))


def _get_text_area_for_images(images: list[ImageData]):
    """이미지 위치에 따라 텍스트 영역을 조정"""
    if not images:
        return Inches(0.5), Inches(8.6)

    positions = [img.position for img in images]

    if any(p in ("full", "center") for p in positions):
        return None, None

    if any(p.startswith("right") for p in positions):
        return Inches(0.5), Inches(4.8)

    if any(p.startswith("left") for p in positions):
        return Inches(5.0), Inches(4.5)

    return Inches(0.5), Inches(8.6)


# ── Slide Builders ──

def _build_title_slide(prs, data: SlideRequest):
    """표지: 원본 슬라이드1 포맷. placeholder를 직접 사용."""
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_TITLE])
    _set_slide_bg(slide, COVER_BG)

    for ph in slide.placeholders:
        idx = ph.placeholder_format.idx
        if idx == 0:  # CENTER_TITLE
            _set_placeholder_text(ph, data.title, font_size=Pt(40),
                                  font_color=RGBColor(0xFF, 0xFF, 0xFF),
                                  alignment=PP_ALIGN.CENTER)
        elif idx == 1:  # SUBTITLE
            _set_placeholder_text(ph, data.subtitle or "",
                                  font_size=Pt(14),
                                  font_color=RGBColor(0xBB, 0xBB, 0xBB),
                                  alignment=PP_ALIGN.RIGHT)
        else:
            _remove_placeholder(slide, idx)


def _build_toc_slide(prs, data: SlideRequest):
    """목차: placeholder[0]에 제목, placeholder[1]에 목차 내용."""
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_CONTENT])

    for ph in slide.placeholders:
        idx = ph.placeholder_format.idx
        if idx == 0:
            _set_placeholder_text(ph, "목차", bold=True)
        elif idx == 1:
            if data.bullets:
                tf = ph.text_frame
                tf.clear()
                for i, b in enumerate(data.bullets):
                    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                    if i > 0:
                        p.space_before = Pt(10)
                    p.text = f"{i+1}. {b}"
                    p.font.name = "맑은 고딕"
                    p.font.size = Pt(15)
                    p.font.bold = True
            else:
                _remove_placeholder(slide, idx)


def _build_content_slide(prs, data: SlideRequest):
    """콘텐츠: placeholder 사용 + 이미지는 위치 프리셋에 따라 배치."""
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_CONTENT])

    text_left, text_width = _get_text_area_for_images(data.images)

    for ph in slide.placeholders:
        idx = ph.placeholder_format.idx
        if idx == 0:
            _set_placeholder_text(ph, data.title)
        elif idx == 1:
            if data.bullets and text_left is not None:
                _fill_bullets_in_placeholder(ph, data.bullets)
                ph.left = int(text_left)
                ph.width = int(text_width)
            else:
                _remove_placeholder(slide, idx)

    for img in data.images:
        _add_positioned_image(slide, img)


def _build_image_text_slide(prs, data: SlideRequest):
    """이미지+텍스트: 이미지 위치 프리셋 지원."""
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_CONTENT])

    for ph in slide.placeholders:
        idx = ph.placeholder_format.idx
        if idx == 0:
            _set_placeholder_text(ph, data.title)
        elif idx == 1:
            _remove_placeholder(slide, idx)

    for img in data.images:
        _add_positioned_image(slide, img)

    caption = data.caption or (data.bullets[0] if data.bullets else "")
    if caption:
        _add_textbox(
            slide, Inches(0.5), Inches(4.4), Inches(9), Inches(0.4),
            caption, font_size=Pt(15), bold=True,
        )


def _build_grid_slide(prs, data: SlideRequest):
    """이미지 그리드."""
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_CONTENT])

    for ph in slide.placeholders:
        idx = ph.placeholder_format.idx
        if idx == 0:
            _set_placeholder_text(ph, data.title)
        elif idx == 1:
            _remove_placeholder(slide, idx)

    images = data.images
    labels = data.bullets
    count = len(images)
    if count == 0:
        return

    cols = min(count, 3)
    card_w = Inches(2.5)
    card_h = Inches(1.6)
    gap = Inches(0.3)
    total_w = cols * card_w + (cols - 1) * gap
    start_x = (Inches(10) - total_w) // 2

    for i, img in enumerate(images):
        col = i % cols
        row = i // cols
        if row >= 2:
            break
        x = start_x + col * (card_w + gap)
        y = Inches(1.2) + row * (card_h + Inches(0.6))

        _add_image(slide, img, x, y, card_w, card_h)

        label = labels[i] if i < len(labels) else ""
        if label:
            _add_textbox(
                slide, x, y + card_h + Inches(0.05), card_w, Inches(0.25),
                label, font_size=Pt(10), alignment=PP_ALIGN.CENTER,
            )


def _build_two_column_slide(prs, data: SlideRequest):
    """2단 비교: 비교 레이아웃의 placeholder 직접 사용."""
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_COMPARE])

    for ph in slide.placeholders:
        idx = ph.placeholder_format.idx
        if idx == 0:
            _set_placeholder_text(ph, data.title)
        elif idx == 1:
            _set_placeholder_text(ph, data.left_title, bold=True)
        elif idx == 2:
            _fill_bullets_in_placeholder(ph, data.left_bullets)
        elif idx == 3:
            _set_placeholder_text(ph, data.right_title, bold=True)
        elif idx == 4:
            _fill_bullets_in_placeholder(ph, data.right_bullets)


def _build_closing_slide(prs, data: SlideRequest):
    """마무리: 원본 슬라이드13 포맷."""
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_CONTENT])

    for ph in slide.placeholders:
        idx = ph.placeholder_format.idx
        if idx == 0:
            _set_placeholder_text(ph, "맺음말")
        elif idx == 1:
            _remove_placeholder(slide, idx)

    _add_textbox(
        slide, Inches(2.80), Inches(2.16), Inches(7.00), Inches(2.04),
        data.title or "감사합니다.",
        font_size=Pt(54), alignment=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE,
    )


LAYOUT_BUILDERS = {
    "title": _build_title_slide,
    "toc": _build_toc_slide,
    "content": _build_content_slide,
    "image-text": _build_image_text_slide,
    "grid": _build_grid_slide,
    "two-column": _build_two_column_slide,
    "closing": _build_closing_slide,
}


def generate_pptx(title: str, theme_name: str, slides: list[SlideRequest]) -> BytesIO:
    prs = Presentation(TEMPLATE_PATH)

    for slide_data in slides:
        builder = LAYOUT_BUILDERS.get(slide_data.layout)
        if builder:
            builder(prs, slide_data)

    buffer = BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer
